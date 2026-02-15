#!/usr/bin/env python3
"""
Extract data from Week 33 reporting files
"""

import openpyxl
from pptx import Presentation
import json
import sys
from pathlib import Path

def extract_pptx_data(file_path):
    """Extract text content from PowerPoint file"""
    print(f"\n{'='*80}")
    print(f"EXTRACTING: {file_path}")
    print(f"{'='*80}")

    prs = Presentation(file_path)
    data = {
        "highlights": [],
        "lowlights": [],
        "emerging_issues": [],
        "priorities": [],
        "all_text": []
    }

    current_section = None

    for slide_num, slide in enumerate(prs.slides, 1):
        print(f"\n--- Slide {slide_num} ---")
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                print(text)
                data["all_text"].append(text)

                # Try to categorize content
                text_lower = text.lower()
                if "highlight" in text_lower:
                    current_section = "highlights"
                elif "lowlight" in text_lower:
                    current_section = "lowlights"
                elif "emerging" in text_lower or "issue" in text_lower:
                    current_section = "emerging_issues"
                elif "priorit" in text_lower or "action" in text_lower:
                    current_section = "priorities"
                elif current_section and text not in ["Highlights", "Lowlights", "Emerging Issues", "Priorities", "Actions"]:
                    data[current_section].append(text)

    return data

def extract_excel_data(file_path):
    """Extract data from Excel file"""
    print(f"\n{'='*80}")
    print(f"EXTRACTING: {file_path}")
    print(f"{'='*80}")

    wb = openpyxl.load_workbook(file_path, data_only=True)

    data = {
        "file_name": Path(file_path).name,
        "sheets": {},
        "equipment_availability": {},
        "breakdowns": [],
        "safety_incidents": [],
        "service_compliance": {},
        "weekly_averages": {}
    }

    print(f"\nAvailable sheets: {wb.sheetnames}")

    for sheet_name in wb.sheetnames:
        print(f"\n--- Sheet: {sheet_name} ---")
        ws = wb[sheet_name]

        sheet_data = []
        max_rows = min(ws.max_row, 100)  # Limit to first 100 rows for initial inspection

        for row_idx, row in enumerate(ws.iter_rows(min_row=1, max_row=max_rows, values_only=True), 1):
            # Filter out completely empty rows
            if any(cell is not None and str(cell).strip() != '' for cell in row):
                # Convert row to list and clean up
                cleaned_row = []
                for cell in row:
                    if cell is None:
                        cleaned_row.append(None)
                    elif isinstance(cell, (int, float)):
                        cleaned_row.append(cell)
                    else:
                        cleaned_row.append(str(cell).strip())

                sheet_data.append({
                    "row_num": row_idx,
                    "data": cleaned_row
                })

                # Print first 30 rows for inspection
                if row_idx <= 30:
                    print(f"Row {row_idx}: {cleaned_row[:15]}")  # Show first 15 columns

        data["sheets"][sheet_name] = sheet_data

        # Try to find specific data patterns
        for row_info in sheet_data:
            row = row_info["data"]
            row_text = ' '.join([str(cell) for cell in row if cell is not None]).lower()

            # Look for equipment types
            for equipment_type in ['dt', 'fl', 'hd', 'rt', 'sr', 'uv']:
                if equipment_type in row_text or equipment_type.upper() in row_text:
                    # Look for percentage values in the row
                    for cell in row:
                        if isinstance(cell, (int, float)) and 0 <= cell <= 100:
                            if equipment_type.upper() not in data["equipment_availability"]:
                                data["equipment_availability"][equipment_type.upper()] = []
                            data["equipment_availability"][equipment_type.upper()].append({
                                "value": cell,
                                "sheet": sheet_name,
                                "row": row_info["row_num"],
                                "context": row[:10]
                            })

            # Look for breakdown-related data
            if any(keyword in row_text for keyword in ['breakdown', 'failure', 'fault', 'issue']):
                data["breakdowns"].append({
                    "sheet": sheet_name,
                    "row": row_info["row_num"],
                    "content": row[:15]
                })

            # Look for safety incidents
            if any(keyword in row_text for keyword in ['safety', 'incident', 'accident', 'injury']):
                data["safety_incidents"].append({
                    "sheet": sheet_name,
                    "row": row_info["row_num"],
                    "content": row[:15]
                })

            # Look for availability or average percentages
            if 'availab' in row_text or 'average' in row_text or 'week' in row_text:
                for cell in row:
                    if isinstance(cell, (int, float)) and 0 <= cell <= 100:
                        data["weekly_averages"][f"{sheet_name}_row_{row_info['row_num']}"] = {
                            "value": cell,
                            "context": row[:10]
                        }

    return data

def main():
    base_path = Path("/home/user/weekly-report-generator/data/week33")

    files_to_process = [
        ("Gloria HEAL Week 33.pptx", "pptx"),
        ("N3 Eng Report Week  06 Feb - 12 Feb 2026.xlsx", "xlsx"),
        ("Nch2 Weekly Report 13 February 2026.xlsx", "xlsx"),
        ("Gloria Eng Report Week 06-12 Feb 2026.xlsx", "xlsx")
    ]

    all_data = {}

    for filename, file_type in files_to_process:
        file_path = base_path / filename

        try:
            if file_type == "pptx":
                all_data[filename] = extract_pptx_data(str(file_path))
            elif file_type == "xlsx":
                all_data[filename] = extract_excel_data(str(file_path))
        except Exception as e:
            print(f"\nERROR processing {filename}: {str(e)}")
            import traceback
            traceback.print_exc()

    # Save extracted data to JSON
    output_file = base_path / "extracted_data.json"
    with open(output_file, 'w') as f:
        json.dump(all_data, f, indent=2, default=str)

    print(f"\n{'='*80}")
    print(f"Extraction complete! Data saved to: {output_file}")
    print(f"{'='*80}")

    return all_data

if __name__ == "__main__":
    main()

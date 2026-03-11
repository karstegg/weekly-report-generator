"""
HEAL Data Extractor

Extracts Highlights, Emerging Issues, Actions/Lowlights, Lookahead/Priorities
from:
  - Text files (e.g., "N3 HEAL Page Week31.txt")
  - PowerPoint files (e.g., "HEAL page (002).pptx")

Text file format:
    Highlights
    * Topic: Description
    * Topic: Description

    Lowlights
    * Topic: Description
    ...
"""
from pathlib import Path
import re

try:
    from pptx import Presentation
    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


def extract_heal_from_text(filepath: Path, site: str) -> dict:
    """Parse a HEAL text file into structured data."""
    text = filepath.read_text(encoding="utf-8")
    sections = {
        "highlights": [],
        "lowlights": [],
        "emergingIssues": [],
        "priorities": [],
    }

    # Map section headers to keys
    header_map = {
        "highlights": "highlights",
        "lowlights": "lowlights",
        "emerging issues": "emergingIssues",
        "emerging": "emergingIssues",
        "priorities": "priorities",
        "actions": "lowlights",
        "lookahead": "priorities",
    }

    current_section = None
    for line in text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue

        # Check if this is a section header
        lower = stripped.lower()
        if lower in header_map:
            current_section = header_map[lower]
            continue

        # Check if this is a bullet point
        if current_section and (stripped.startswith("*") or stripped.startswith("-") or stripped.startswith("•")):
            bullet_text = re.sub(r"^[*\-•]\s*", "", stripped)
            sections[current_section].append({
                "site": site,
                "text": bullet_text,
            })

    return sections


def extract_heal_from_pptx(filepath: Path, site: str) -> dict:
    """Extract HEAL data from a PowerPoint file."""
    if not HAS_PPTX:
        print(f"  python-pptx not installed, skipping {filepath.name}")
        return {}

    prs = Presentation(str(filepath))
    sections = {
        "highlights": [],
        "lowlights": [],
        "emergingIssues": [],
        "priorities": [],
    }

    # Extract all text from all slides
    all_text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        all_text.append(text)
            if shape.has_table:
                for row in shape.table.rows:
                    for cell in row.cells:
                        text = cell.text.strip()
                        if text:
                            all_text.append(text)

    # Parse the extracted text as if it were a text file
    header_map = {
        "highlights": "highlights",
        "lowlights": "lowlights",
        "emerging issues": "emergingIssues",
        "priorities": "priorities",
    }

    current_section = None
    for text in all_text:
        lower = text.lower()
        if lower in header_map:
            current_section = header_map[lower]
            continue

        if current_section and len(text) > 3:
            bullet_text = re.sub(r"^[*\-•]\s*", "", text)
            sections[current_section].append({
                "site": site,
                "text": bullet_text,
            })

    return sections


def extract_heal_data(data_dir: Path, week: int) -> dict:
    """
    Extract HEAL data from all available sources.
    Returns combined HEAL data for all sites.
    """
    combined = {
        "highlights": [],
        "lowlights": [],
        "emergingIssues": [],
        "priorities": [],
    }

    site_map = {
        "N3": ["N3", "Nchwaning 3", "Nchwaning3"],
        "N2": ["N2", "Nchwaning 2", "Nchwaning2"],
        "Gloria": ["Gloria"],
    }

    for site, prefixes in site_map.items():
        found = False

        # Try text files first
        for prefix in prefixes:
            patterns = [
                f"{prefix} HEAL Page Week{week}.txt",
                f"{prefix} HEAL Page - Week{week}.txt",
                f"{prefix}_HEAL_Week{week}.txt",
            ]
            for pattern in patterns:
                filepath = data_dir / pattern
                if filepath.exists():
                    print(f"  Found text HEAL: {filepath.name}")
                    sections = extract_heal_from_text(filepath, site)
                    for key in combined:
                        combined[key].extend(sections.get(key, []))
                    found = True
                    break
            if found:
                break

        # Try PPTX fallback
        if not found:
            for f in data_dir.glob("*.pptx"):
                if any(p.lower() in f.name.lower() for p in prefixes) and "heal" in f.name.lower():
                    print(f"  Found PPTX HEAL: {f.name}")
                    sections = extract_heal_from_pptx(f, site)
                    for key in combined:
                        combined[key].extend(sections.get(key, []))
                    found = True
                    break

        if not found:
            print(f"  No HEAL data found for {site}")

    return combined

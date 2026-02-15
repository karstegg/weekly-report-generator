#!/usr/bin/env python3
"""
Detailed extraction of Week 33 data for weekly reporting
"""

import openpyxl
from pptx import Presentation
import json
from pathlib import Path
from typing import Dict, List, Any

def extract_n3_engineering_report(file_path: str) -> Dict[str, Any]:
    """Extract detailed data from N3 Engineering Report"""
    print(f"\n{'='*80}")
    print(f"N3 ENGINEERING REPORT")
    print(f"{'='*80}")

    wb = openpyxl.load_workbook(file_path, data_only=True)
    data = {
        "equipment_availability": {},
        "safety": {},
        "service_compliance": {},
        "breakdowns": [],
        "hydraulic_hoses": {},
        "oil_consumption": {}
    }

    # Extract from Weekly report sheet
    if "Weekly report" in wb.sheetnames:
        ws = wb["Weekly report"]
        print("\n--- SAFETY DATA ---")

        # Safety data (rows 8-11)
        data["safety"] = {
            "minor_treatment_case": {"ytd": ws.cell(8, 2).value, "week": ws.cell(8, 3).value},
            "lost_time_injury": {"ytd": ws.cell(9, 2).value, "week": ws.cell(9, 3).value},
            "property_damage": {"ytd": ws.cell(10, 2).value, "week": ws.cell(10, 3).value},
            "incidents_reported": {"ytd": ws.cell(11, 2).value, "week": ws.cell(11, 3).value}
        }

        print(f"MTC: YTD={data['safety']['minor_treatment_case']['ytd']}, Week={data['safety']['minor_treatment_case']['week']}")
        print(f"LTI: YTD={data['safety']['lost_time_injury']['ytd']}, Week={data['safety']['lost_time_injury']['week']}")
        print(f"PD: YTD={data['safety']['property_damage']['ytd']}, Week={data['safety']['property_damage']['week']}")

        # Look for equipment availability data - scan the sheet
        print("\n--- EQUIPMENT AVAILABILITY ---")
        for row_idx in range(1, min(ws.max_row, 150)):
            row_values = [ws.cell(row_idx, col).value for col in range(1, 20)]
            row_text = ' '.join([str(v) for v in row_values if v is not None]).lower()

            # Look for availability percentages by equipment type
            if any(keyword in row_text for keyword in ['availability', 'avail %', 'weekly average']):
                # Check if equipment type is mentioned
                for col_idx in range(1, 20):
                    cell_val = ws.cell(row_idx, col_idx).value
                    if cell_val and isinstance(cell_val, str):
                        cell_lower = cell_val.lower()
                        if 'dt' in cell_lower or 'dump truck' in cell_lower:
                            # Look for percentage in nearby cells
                            for offset in range(1, 10):
                                val = ws.cell(row_idx, col_idx + offset).value
                                if isinstance(val, (int, float)) and 0 <= val <= 1:
                                    data['equipment_availability']['DT'] = val * 100
                                    print(f"DT Availability: {val * 100:.2f}%")
                                    break
                        elif 'fl' in cell_lower or 'loader' in cell_lower:
                            for offset in range(1, 10):
                                val = ws.cell(row_idx, col_idx + offset).value
                                if isinstance(val, (int, float)) and 0 <= val <= 1:
                                    data['equipment_availability']['FL'] = val * 100
                                    print(f"FL Availability: {val * 100:.2f}%")
                                    break
                        elif 'hd' in cell_lower or 'drill' in cell_lower:
                            for offset in range(1, 10):
                                val = ws.cell(row_idx, col_idx + offset).value
                                if isinstance(val, (int, float)) and 0 <= val <= 1:
                                    data['equipment_availability']['HD'] = val * 100
                                    print(f"HD Availability: {val * 100:.2f}%")
                                    break
                        elif 'rt' in cell_lower or 'bolter' in cell_lower:
                            for offset in range(1, 10):
                                val = ws.cell(row_idx, col_idx + offset).value
                                if isinstance(val, (int, float)) and 0 <= val <= 1:
                                    data['equipment_availability']['RT'] = val * 100
                                    print(f"RT Availability: {val * 100:.2f}%")
                                    break
                        elif 'sr' in cell_lower or 'scaler' in cell_lower:
                            for offset in range(1, 10):
                                val = ws.cell(row_idx, col_idx + offset).value
                                if isinstance(val, (int, float)) and 0 <= val <= 1:
                                    data['equipment_availability']['SR'] = val * 100
                                    print(f"SR Availability: {val * 100:.2f}%")
                                    break

    # Extract from Compliance Report sheet
    if "Compliance Report" in wb.sheetnames:
        ws = wb["Compliance Report"]
        print("\n--- SERVICE COMPLIANCE ---")

        for row_idx in range(6, 20):
            fleet_type = ws.cell(row_idx, 2).value
            if fleet_type:
                planned = ws.cell(row_idx, 4).value
                actual = ws.cell(row_idx, 6).value
                compliance = ws.cell(row_idx, 8).value
                comment = ws.cell(row_idx, 12).value

                if fleet_type and (planned or actual):
                    data['service_compliance'][str(fleet_type)] = {
                        "planned": planned,
                        "actual": actual,
                        "compliance": compliance,
                        "comment": comment
                    }
                    print(f"{fleet_type}: Planned={planned}, Actual={actual}, Compliance={compliance}")

    # Extract hydraulic hose data
    if "Hydraulic pipe report" in wb.sheetnames:
        ws = wb["Hydraulic pipe report"]
        print("\n--- HYDRAULIC HOSE REPORT ---")

        for row_idx in range(4, 12):
            equipment = ws.cell(row_idx, 1).value
            if equipment:
                data['hydraulic_hoses'][str(equipment)] = {
                    "abrasion": ws.cell(row_idx, 2).value or 0,
                    "fitting_damaged": ws.cell(row_idx, 3).value or 0,
                    "hose_burst": ws.cell(row_idx, 4).value or 0,
                    "hose_damaged": ws.cell(row_idx, 5).value or 0,
                    "service": ws.cell(row_idx, 9).value or 0,
                    "total": ws.cell(row_idx, 10).value or 0
                }
                print(f"{equipment}: {data['hydraulic_hoses'][str(equipment)]}")

    return data

def extract_nch2_report(file_path: str) -> Dict[str, Any]:
    """Extract detailed data from Nch2 Weekly Report"""
    print(f"\n{'='*80}")
    print(f"NCH2 WEEKLY REPORT")
    print(f"{'='*80}")

    wb = openpyxl.load_workbook(file_path, data_only=True)
    data = {
        "equipment_availability": {},
        "safety": {},
        "service_compliance": {},
        "production": {}
    }

    # Extract from Eng Weekly report sheet
    if "Eng Weekly report" in wb.sheetnames:
        ws = wb["Eng Weekly report"]

        print("\n--- PRODUCTION DATA ---")
        data["production"] = {
            "rom_plan_ytd": ws.cell(3, 8).value,
            "rom_actual_ytd": ws.cell(3, 9).value,
            "rom_plan_mtd": ws.cell(3, 10).value,
            "rom_actual_mtd": ws.cell(3, 11).value,
            "product_plan_ytd": ws.cell(4, 8).value,
            "product_actual_ytd": ws.cell(4, 9).value,
            "product_plan_mtd": ws.cell(4, 10).value,
            "product_actual_mtd": ws.cell(4, 11).value
        }

        print(f"ROM: MTD Plan={data['production']['rom_plan_mtd']}, Actual={data['production']['rom_actual_mtd']}")
        print(f"Product: MTD Plan={data['production']['product_plan_mtd']}, Actual={data['production']['product_actual_mtd']}")

        print("\n--- SAFETY DATA ---")
        data["safety"] = {
            "mtc": {"ytd": ws.cell(8, 2).value, "week": ws.cell(8, 3).value},
            "lti": {"ytd": ws.cell(9, 2).value, "week": ws.cell(9, 3).value},
            "pd": {"ytd": ws.cell(10, 2).value, "week": ws.cell(10, 3).value},
            "incidents": {"ytd": ws.cell(11, 2).value, "week": ws.cell(11, 3).value}
        }

        print(f"MTC: YTD={data['safety']['mtc']['ytd']}, Week={data['safety']['mtc']['week']}")
        print(f"LTI: YTD={data['safety']['lti']['ytd']}, Week={data['safety']['lti']['week']}")
        print(f"PD: YTD={data['safety']['pd']['ytd']}, Week={data['safety']['pd']['week']}")

    # Extract from Compliance Report sheet
    if "Compliance Report" in wb.sheetnames:
        ws = wb["Compliance Report"]
        print("\n--- SERVICE COMPLIANCE ---")

        for row_idx in range(7, 15):
            fleet_type = ws.cell(row_idx, 2).value
            if fleet_type:
                planned = ws.cell(row_idx, 3).value
                actual = ws.cell(row_idx, 4).value
                compliance = ws.cell(row_idx, 5).value
                comment = ws.cell(row_idx, 12).value

                if planned or actual:
                    data['service_compliance'][str(fleet_type)] = {
                        "planned": planned,
                        "actual": actual,
                        "compliance": compliance,
                        "comment": comment
                    }
                    print(f"{fleet_type}: Planned={planned}, Actual={actual}, Compliance={compliance}")

    # Extract Fermel availability data
    if "Fermel Weekly" in wb.sheetnames:
        ws = wb["Fermel Weekly"]
        print("\n--- FERMEL AVAILABILITY ---")

        # Look for total availability (usually in rows 3-15, columns 10-12)
        for row_idx in range(3, 16):
            description = ws.cell(row_idx, 9).value
            total = ws.cell(row_idx, 10).value
            breakdown = ws.cell(row_idx, 11).value
            availability = ws.cell(row_idx, 12).value

            if description and availability is not None:
                print(f"{description}: Total={total}, Breakdown={breakdown}, Availability={availability}")
                if isinstance(availability, (int, float)):
                    data['equipment_availability'][str(description)] = {
                        "total": total,
                        "breakdown": breakdown,
                        "availability": availability * 100 if availability <= 1 else availability
                    }

    return data

def extract_gloria_report(file_path: str) -> Dict[str, Any]:
    """Extract detailed data from Gloria Engineering Report"""
    print(f"\n{'='*80}")
    print(f"GLORIA ENGINEERING REPORT")
    print(f"{'='*80}")

    wb = openpyxl.load_workbook(file_path, data_only=True)
    data = {
        "equipment_availability": {},
        "safety": {},
        "production": {},
        "service_compliance": {}
    }

    # Extract from Weekly report sheet
    if "Weekly report" in wb.sheetnames:
        ws = wb["Weekly report"]

        print("\n--- PRODUCTION DATA ---")
        # Production is usually in early rows
        for row_idx in range(1, 10):
            for col_idx in range(1, 15):
                cell_val = ws.cell(row_idx, col_idx).value
                if cell_val and isinstance(cell_val, str):
                    if 'rom' in cell_val.lower():
                        # Look for values in nearby cells
                        plan_ytd = ws.cell(row_idx, col_idx + 1).value
                        actual_ytd = ws.cell(row_idx, col_idx + 2).value
                        if plan_ytd or actual_ytd:
                            data['production']['rom_plan_ytd'] = plan_ytd
                            data['production']['rom_actual_ytd'] = actual_ytd
                            print(f"ROM: Plan YTD={plan_ytd}, Actual YTD={actual_ytd}")
                    elif 'product' in cell_val.lower() and 'rom' not in cell_val.lower():
                        plan_ytd = ws.cell(row_idx, col_idx + 1).value
                        actual_ytd = ws.cell(row_idx, col_idx + 2).value
                        if plan_ytd or actual_ytd:
                            data['production']['product_plan_ytd'] = plan_ytd
                            data['production']['product_actual_ytd'] = actual_ytd
                            print(f"Product: Plan YTD={plan_ytd}, Actual YTD={actual_ytd}")

        print("\n--- SAFETY DATA ---")
        # Safety data - scan for safety section
        for row_idx in range(1, 50):
            cell_val = ws.cell(row_idx, 1).value
            if cell_val and isinstance(cell_val, str):
                if 'minor treatment' in cell_val.lower() or 'mtc' in cell_val.lower():
                    data['safety']['mtc'] = {
                        "ytd": ws.cell(row_idx, 2).value,
                        "week": ws.cell(row_idx, 3).value
                    }
                elif 'lost time' in cell_val.lower() or cell_val.lower() == 'lti':
                    data['safety']['lti'] = {
                        "ytd": ws.cell(row_idx, 2).value,
                        "week": ws.cell(row_idx, 3).value
                    }
                elif 'property damage' in cell_val.lower() or cell_val.lower() == 'pd':
                    data['safety']['pd'] = {
                        "ytd": ws.cell(row_idx, 2).value,
                        "week": ws.cell(row_idx, 3).value
                    }

        if data['safety']:
            for key, val in data['safety'].items():
                print(f"{key.upper()}: YTD={val.get('ytd')}, Week={val.get('week')}")

        print("\n--- EQUIPMENT AVAILABILITY ---")
        # Scan for availability data
        for row_idx in range(1, 150):
            for col_idx in range(1, 15):
                cell_val = ws.cell(row_idx, col_idx).value
                if cell_val and isinstance(cell_val, str):
                    cell_lower = cell_val.lower()

                    # Look for equipment types with availability
                    equipment_type = None
                    if 'dump truck' in cell_lower or (cell_lower.strip() == 'dt' or 'dt ' in cell_lower):
                        equipment_type = 'DT'
                    elif 'loader' in cell_lower or cell_lower.strip() == 'fl':
                        equipment_type = 'FL'
                    elif 'drill' in cell_lower or cell_lower.strip() == 'hd':
                        equipment_type = 'HD'
                    elif 'bolter' in cell_lower or cell_lower.strip() == 'rt':
                        equipment_type = 'RT'
                    elif 'scaler' in cell_lower or cell_lower.strip() == 'sr':
                        equipment_type = 'SR'

                    if equipment_type and equipment_type not in data['equipment_availability']:
                        # Look for percentage in nearby cells
                        for offset in range(1, 12):
                            val = ws.cell(row_idx, col_idx + offset).value
                            if isinstance(val, (int, float)) and 0 <= val <= 1:
                                data['equipment_availability'][equipment_type] = val * 100
                                print(f"{equipment_type}: {val * 100:.2f}%")
                                break

    # Extract from Compliance Report if exists
    for sheet_name in wb.sheetnames:
        if 'compliance' in sheet_name.lower():
            ws = wb[sheet_name]
            print(f"\n--- SERVICE COMPLIANCE ({sheet_name}) ---")

            for row_idx in range(1, 30):
                fleet_type = ws.cell(row_idx, 2).value
                if fleet_type and isinstance(fleet_type, str):
                    planned = ws.cell(row_idx, 4).value
                    actual = ws.cell(row_idx, 6).value
                    compliance = ws.cell(row_idx, 8).value

                    if (planned or actual) and fleet_type.strip():
                        data['service_compliance'][fleet_type.strip()] = {
                            "planned": planned,
                            "actual": actual,
                            "compliance": compliance
                        }
                        print(f"{fleet_type}: Planned={planned}, Actual={actual}, Compliance={compliance}")
            break

    return data

def extract_heal_pptx(file_path: str) -> Dict[str, Any]:
    """Extract HEAL data from PowerPoint"""
    print(f"\n{'='*80}")
    print(f"GLORIA HEAL PRESENTATION")
    print(f"{'='*80}")

    prs = Presentation(file_path)
    data = {
        "highlights": [],
        "lowlights": [],
        "emerging_issues": [],
        "priorities": []
    }

    current_section = None

    for slide_num, slide in enumerate(prs.slides, 1):
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                text = shape.text.strip()
                text_lower = text.lower()

                # Identify sections
                if "highlight" in text_lower and len(text) < 50:
                    current_section = "highlights"
                    print(f"\n--- HIGHLIGHTS ---")
                elif "lowlight" in text_lower and len(text) < 50:
                    current_section = "lowlights"
                    print(f"\n--- LOWLIGHTS ---")
                elif "emerging" in text_lower and len(text) < 50:
                    current_section = "emerging_issues"
                    print(f"\n--- EMERGING ISSUES ---")
                elif "priorit" in text_lower or ("action" in text_lower and len(text) < 50):
                    current_section = "priorities"
                    print(f"\n--- PRIORITIES/ACTIONS ---")
                elif current_section and len(text) > 10:
                    # Add content to current section
                    if text not in ["Highlights", "Lowlights", "Emerging Issues", "Priorities", "Actions", "High & Low"]:
                        data[current_section].append(text)
                        print(f"  • {text}")

    return data

def main():
    base_path = Path("/home/user/weekly-report-generator/data/week33")

    # Also read the TXT files for HEAL data
    heal_data = {}

    n2_heal = base_path / "N2 HEAL PAge - Week 33.txt"
    if n2_heal.exists():
        with open(n2_heal, 'r') as f:
            heal_data['n2_heal'] = f.read()

    n3_heal = base_path / "N3 HEAL PAge - Week 33.txt"
    if n3_heal.exists():
        with open(n3_heal, 'r') as f:
            heal_data['n3_heal'] = f.read()

    # Extract all data
    week33_data = {
        "heal": {
            "gloria_pptx": extract_heal_pptx(str(base_path / "Gloria HEAL Week 33.pptx")),
            "n2_text": heal_data.get('n2_heal', ''),
            "n3_text": heal_data.get('n3_heal', '')
        },
        "n3_engineering": extract_n3_engineering_report(str(base_path / "N3 Eng Report Week  06 Feb - 12 Feb 2026.xlsx")),
        "nch2_weekly": extract_nch2_report(str(base_path / "Nch2 Weekly Report 13 February 2026.xlsx")),
        "gloria_engineering": extract_gloria_report(str(base_path / "Gloria Eng Report Week 06-12 Feb 2026.xlsx"))
    }

    # Save to JSON
    output_file = base_path / "week33_structured_data.json"
    with open(output_file, 'w') as f:
        json.dump(week33_data, f, indent=2, default=str)

    print(f"\n{'='*80}")
    print(f"Data saved to: {output_file}")
    print(f"{'='*80}")

    return week33_data

if __name__ == "__main__":
    main()

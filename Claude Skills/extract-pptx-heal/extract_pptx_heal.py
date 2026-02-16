#!/usr/bin/env python3
"""
Extract HEAL content from PowerPoint files and save as formatted text.

Supports multiple table layouts:
- Gloria: 4x2 table with Highlights/Lowlights in left/right columns
- Shafts & Winders: Multiple single-cell tables with section headers
- N3: Detected from file structure
"""

import argparse
import sys
import os
from pathlib import Path
from pptx import Presentation

# Configure UTF-8 encoding for console output (prevent charmap errors)
if hasattr(sys.stdout, 'reconfigure'):
    sys.stdout.reconfigure(encoding='utf-8')
if hasattr(sys.stderr, 'reconfigure'):
    sys.stderr.reconfigure(encoding='utf-8')


def find_weekly_reports_root():
    """Auto-detect Weekly Reports root directory from any execution location"""
    current = os.getcwd()

    # Check if we're already in Weekly Reports directory
    if os.path.basename(current) == "Weekly Reports":
        return current

    # Check if we're in a skill subdirectory
    if ".claude" in current and "skills" in current:
        # Go up to Weekly Reports root
        parts = current.split(os.sep)
        try:
            idx = parts.index("Weekly Reports")
            return os.sep.join(parts[:idx+1])
        except ValueError:
            pass

    # Last resort: search parent directories for Week 1 folder
    test_path = current
    for _ in range(5):  # Check up to 5 levels up
        if os.path.exists(os.path.join(test_path, "Week 1")):
            return test_path
        test_path = os.path.dirname(test_path)

    return None


# Change to correct working directory at script start
weekly_reports_root = find_weekly_reports_root()
if weekly_reports_root:
    os.chdir(weekly_reports_root)
else:
    print("ERROR: Could not find Weekly Reports root directory", file=sys.stderr)
    print("This script must be run from within the Weekly Reports directory tree", file=sys.stderr)
    sys.exit(1)


def find_pptx_file(week_num, site):
    """Find the PPTX file in Week directory matching site pattern."""
    week_dir = Path(f"Week {week_num}")

    if not week_dir.exists():
        return None

    # Site-specific search patterns
    patterns = {
        "gloria": ["HEAL page*.pptx", "*gloria*.pptx"],
        "shafts-winders": ["*SHAFTS*WINDERS*.pptx", "*S&W*.pptx"],
        "n3": ["*Heal*N3*.pptx", "*N3*HEAL*.pptx", "*Nchwaning 3*HEAL*.pptx"],
    }

    for pattern in patterns.get(site, []):
        for pptx_file in week_dir.glob(pattern):
            return pptx_file

    return None


def extract_gloria_heal(pptx_path):
    """Extract HEAL content from Gloria's 4x2 table layout."""
    prs = Presentation(pptx_path)
    heal_data = {
        "Highlights": [],
        "Lowlights": [],
        "Emerging Issues": [],
        "Priorities": [],
    }

    for slide in prs.slides:
        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table

                # Gloria layout: 4 rows x 2 columns
                if len(table.rows) >= 4 and len(table.columns) >= 2:
                    # Row 1: Highlights (left) and Lowlights (right)
                    left_col = table.rows[1].cells[0].text.strip()
                    right_col = table.rows[1].cells[1].text.strip()

                    if left_col:
                        for line in left_col.split('\n'):
                            line = line.strip()
                            if line and line not in heal_data["Highlights"]:
                                heal_data["Highlights"].append(line)

                    if right_col:
                        for line in right_col.split('\n'):
                            line = line.strip()
                            if line and line not in heal_data["Lowlights"]:
                                heal_data["Lowlights"].append(line)

                    # Row 3: Emerging Issues (left) and Priorities (right)
                    if len(table.rows) > 3:
                        left_col = table.rows[3].cells[0].text.strip()
                        right_col = table.rows[3].cells[1].text.strip()

                        if left_col and left_col != "1.":  # Skip numbering artifacts
                            for line in left_col.split('\n'):
                                line = line.strip()
                                if line and line not in heal_data["Emerging Issues"]:
                                    heal_data["Emerging Issues"].append(line)

                        if right_col:
                            for line in right_col.split('\n'):
                                line = line.strip()
                                if line and line not in heal_data["Priorities"]:
                                    heal_data["Priorities"].append(line)

    return heal_data


def extract_shafts_winders_heal(pptx_path):
    """Extract HEAL content from Shafts & Winders' multi-table layout."""
    prs = Presentation(pptx_path)
    heal_data = {
        "Highlights": [],
        "Lowlights": [],
        "Emerging Issues": [],
        "Priorities": [],
        "Actions": [],
    }

    # Process only first slide (contains HEAL content, other slides are schedules)
    if len(prs.slides) > 0:
        slide = prs.slides[0]

        # Each table represents a section
        current_section = None

        for shape in slide.shapes:
            if shape.has_table:
                table = shape.table
                cell_text = table.rows[0].cells[0].text.strip()

                # Identify section from table content
                if cell_text.startswith("Highlights"):
                    current_section = "Highlights"
                elif cell_text.startswith("Lowlights"):
                    current_section = "Lowlights"
                elif cell_text.startswith("Emerging Issues"):
                    current_section = "Emerging Issues"
                elif cell_text.startswith("Priorities"):
                    current_section = "Priorities"
                elif cell_text.startswith("You Can Help By") or cell_text.startswith("Please provide"):
                    current_section = "Actions"

                # Extract items under this section
                if current_section:
                    for line in cell_text.split('\n'):
                        line = line.strip()
                        # Skip section headers and empty lines
                        if line and not line.endswith(":") and "Please provide" not in line:
                            if line not in heal_data[current_section]:
                                heal_data[current_section].append(line)

    return heal_data


def extract_n3_heal(pptx_path):
    """Extract HEAL content from N3 (uses same 4x2 table layout as Gloria)."""
    # N3 uses the same 4x2 table layout as Gloria
    return extract_gloria_heal(pptx_path)


def format_heal_output(heal_data):
    """Format extracted HEAL data into standard text output."""
    import re

    lines = []

    sections = ["Highlights", "Lowlights", "Emerging Issues", "Priorities", "Actions"]

    for section in sections:
        if section in heal_data:
            items = heal_data[section]

            lines.append(section)

            if items:
                for item in items:
                    # Remove numbering artifacts ONLY (e.g., "1.", "2.", "1,", "1)")
                    # Pattern: optional digits followed by punctuation and whitespace at start
                    # This preserves equipment IDs like "74N" while removing "1. " or "2. "
                    item = re.sub(r'^\s*\d+[.,)]\s+', '', item)

                    # Remove leading bullets/dashes
                    item = re.sub(r'^[•–-]\s*', '', item)

                    item = item.strip()

                    if item:
                        lines.append(f"• {item}")

            lines.append("")  # Blank line between sections

    # Remove trailing blank lines
    while lines and lines[-1] == "":
        lines.pop()

    return "\n".join(lines)


def extract_heal(week_num, site):
    """Main extraction function."""

    # Find PPTX file
    pptx_path = find_pptx_file(week_num, site)
    if not pptx_path:
        print(f"Error: No PPTX file found for {site} in Week {week_num}", file=sys.stderr)
        return False

    print(f"Processing: {pptx_path}")

    try:
        # Extract based on site
        if site == "gloria":
            heal_data = extract_gloria_heal(pptx_path)
        elif site == "shafts-winders":
            heal_data = extract_shafts_winders_heal(pptx_path)
        elif site == "n3":
            heal_data = extract_n3_heal(pptx_path)
        else:
            print(f"Error: Unknown site: {site}", file=sys.stderr)
            return False

        # Format output
        output_text = format_heal_output(heal_data)

        # Determine output filename
        site_display = site.replace("-", " & ").title()
        output_filename = f"Week {week_num}/{site_display} HEAL Page Week{week_num}.txt"

        # Save output
        output_path = Path(output_filename)
        output_path.parent.mkdir(parents=True, exist_ok=True)

        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(output_text)

        print("Successfully extracted HEAL content")
        print(f"Output: {output_path}")

        # Show summary
        total_items = sum(len(items) for items in heal_data.values())
        print(f"Extracted {total_items} items across {len([s for s in heal_data.values() if s])} sections")

        return True

    except Exception as e:
        print(f"Error: Failed to extract HEAL content: {e}", file=sys.stderr)
        return False


def main():
    parser = argparse.ArgumentParser(
        description="Extract HEAL content from PowerPoint files and save as formatted text"
    )

    # Support both positional and named arguments
    parser.add_argument(
        "site",
        nargs="?",
        help="Site name: gloria, shafts-winders, or n3"
    )
    parser.add_argument(
        "week",
        nargs="?",
        type=int,
        help="Week number (e.g., 16)"
    )
    parser.add_argument(
        "--site",
        dest="named_site",
        help="Site name (alternative to positional)"
    )
    parser.add_argument(
        "--week",
        type=int,
        dest="named_week",
        help="Week number (alternative to positional)"
    )

    args = parser.parse_args()

    # Determine site and week from arguments
    site = args.site or args.named_site
    week = args.week or args.named_week

    if not site or not week:
        parser.print_help()
        sys.exit(1)

    # Run extraction
    success = extract_heal(week, site)
    sys.exit(0 if success else 1)


if __name__ == "__main__":
    main()
